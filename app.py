import os
# from dotenv import load_dotenv  # Local opcional: no se usa en Streamlit Cloud
import streamlit as st
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
import tempfile
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
)
from pptx import Presentation
from datetime import datetime
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT
from reportlab.lib import colors
from xml.sax.saxutils import escape as xml_escape
# =============== Configuración de credencial ===============
# load_dotenv()  # Si corres local y tienes .env
OPENAI_API_KEY = st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    st.error("Falta OPENAI_API_KEY en tus Secrets (Streamlit Cloud) o variable de entorno.")
    st.stop()

def construir_pdf_conversacion(messages) -> bytes:
    """
    Genera un PDF (bytes) con el historial de chat.
    Usa estilos simples y soporta saltos de línea básicos.
    """
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=36, rightMargin=36,
        topMargin=36, bottomMargin=36
    )

    styles = getSampleStyleSheet()
    title_style = styles["Title"]
    normal = styles["BodyText"]

    # estilo para bloques (un poco más espaciado)
    msg_style = ParagraphStyle(
        "Msg",
        parent=normal,
        fontName="Helvetica",
        fontSize=10.5,
        leading=14,
        alignment=TA_LEFT,
        spaceAfter=8,
        textColor=colors.black
    )
    role_style = ParagraphStyle(
        "Role",
        parent=normal,
        fontName="Helvetica-Bold",
        fontSize=11,
        textColor=colors.HexColor("#0b5394"),
        spaceAfter=2
    )

    story = []
    story.append(Paragraph("CONVERSACIÓN – AGENTE DE BECAS", title_style))
    story.append(Spacer(1, 12))

    for i, msg in enumerate(messages, 1):
        role = "TÚ" if msg["role"] == "user" else "ASISTENTE"
        # escapamos HTML y convertimos saltos de línea a <br/>
        body = xml_escape(msg["content"]).replace("\n", "<br/>")

        story.append(Paragraph(f"{i}. {role}", role_style))
        story.append(Paragraph(body, msg_style))

    doc.build(story)
    pdf_bytes = buffer.getvalue()
    buffer.close()
    return pdf_bytes

# =============== Config página ===============
st.set_page_config(
    page_title="Agente de Becas - Asistente Virtual",
    page_icon="🎓",
    layout="centered"
)

# =============== Estilos ===============
st.markdown("""
<style>
    .chat-container {
        background-color: #f0f2f5;
        padding: 20px;
        border-radius: 10px;
        margin-bottom: 80px;
        max-height: 65vh;
        overflow-y: auto;
    }
    .user-message {
        background: linear-gradient(135deg, #007bff 0%, #0056b3 100%);
        color: white;
        padding: 12px 16px;
        border-radius: 18px 18px 5px 18px;
        margin: 8px 0;
        max-width: 70%;
        margin-left: auto;
        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
        word-wrap: break-word;
    }
    .assistant-message {
        background-color: #ffffff;
        color: #333333;
        padding: 12px 16px;
        border-radius: 18px 18px 18px 5px;
        margin: 8px 0;
        max-width: 70%;
        box-shadow: 0 2px 5px rgba(0,0,0,0.05);
        border: 1px solid #e0e0e0;
        word-wrap: break-word;
    }
    .message-time {
        font-size: 0.7em;
        opacity: 0.7;
        margin-top: 5px;
        text-align: right;
    }
    .assistant-message .message-time {
        text-align: left;
    }
    .control-buttons {
        position: fixed;
        bottom: 80px;
        right: 20px;
        z-index: 1000;
    }
    .control-buttons button {
        margin-left: 5px;
        border-radius: 20px;
        padding: 8px 16px;
    }
    .file-indicator {
        background-color: #e3f2fd;
        border: 1px solid #bbdefb;
        border-radius: 10px;
        padding: 8px 12px;
        margin: 5px 0;
        font-size: 0.9em;
    }
</style>
""", unsafe_allow_html=True)

st.title("Agente de Becas - Asistente Virtual")
st.markdown(
    "¡Hola! Soy tu asistente especializado en becas para estudiantes peruanos. "
    "Pregúntame sobre becas nacionales, internacionales, requisitos, plazos, documentos, etc."
)

# =============== Historial ===============
if "messages" not in st.session_state:
    st.session_state.messages = [
        {"role": "assistant", "content": "¡Hola! ¿En qué puedo ayudarte hoy con respecto a becas?"}
    ]

# =============== Render del chat ===============
st.markdown("<div class='chat-container'>", unsafe_allow_html=True)
for msg in st.session_state.messages:
    timestamp = datetime.now().strftime("%H:%M")
    if msg["role"] == "user":
        st.markdown(f"""
        <div class='user-message'>
            <div>👤 {msg["content"]}</div>
            <div class='message-time'>{timestamp}</div>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown(f"""
        <div class='assistant-message'>
            <div>🎓 {msg["content"]}</div>
            <div class='message-time'>{timestamp}</div>
        </div>
        """, unsafe_allow_html=True)
st.markdown("</div>", unsafe_allow_html=True)

# =============== Botones flotantes ===============
st.markdown("<div class='control-buttons'>", unsafe_allow_html=True)
col1, col2, col3 = st.columns(3)

with col1:
    if st.button("🗑️ Limpiar chat", help="Eliminar toda la conversación"):
        st.session_state.messages = [
            {"role": "assistant", "content": "¡Hola! ¿En qué puedo ayudarte hoy con respecto a becas?"}
        ]
        st.rerun()

with col2:
    if st.button("📥 Exportar chat (TXT)", help="Descargar conversación como texto"):
        if st.session_state.messages:
            export_content = "CONVERSACIÓN - AGENTE DE BECAS\n" + "=" * 50 + "\n\n"
            for i, msg in enumerate(st.session_state.messages, 1):
                role = "TÚ" if msg["role"] == "user" else "ASISTENTE"
                export_content += f"{i}. {role}:\n{msg['content']}\n" + "-" * 30 + "\n\n"
            st.download_button(
                label="💾 Descargar conversación (TXT)",
                data=export_content,
                file_name=f"conversacion_becas_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                mime="text/plain",
                use_container_width=True
            )

with col3:
    if st.button("📄 Exportar PDF", help="Descargar conversación en PDF"):
        if st.session_state.messages:
            pdf_bytes = construir_pdf_conversacion(st.session_state.messages)
            st.download_button(
                label="⬇️ Descargar PDF",
                data=pdf_bytes,
                file_name=f"conversacion_becas_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf",
                mime="application/pdf",
                use_container_width=True
            )
st.markdown("</div>", unsafe_allow_html=True)

# =============== Prompt del sistema ===============
system_prompt = """
Eres un **Agente de Becas Inteligente**, especializado en orientar a estudiantes peruanos sobre oportunidades educativas nacionales e internacionales.

Tu misión:
- Brindar información confiable, clara y actualizada sobre becas de pregrado, posgrado y programas de intercambio.
- Explicar requisitos, plazos, beneficios y procesos de postulación de forma detallada pero sencilla.
- Guiar al usuario con pasos prácticos y consejos estratégicos para aumentar sus posibilidades de éxito.

Tu ámbito de conocimiento incluye:
- **Becas nacionales**: Beca Presidente de la República, Beca 18, Pronabec, Beca Permanencia, Beca Bicentenario entre otras.
- **Becas internacionales**: Fulbright, DAAD, Chevening, Erasmus+, y programas de universidades extranjeras.
- **Aspectos comunes**: requisitos (edad, promedio, idiomas, experiencia), documentos solicitados, cartas de motivación, entrevistas, financiamiento y convenios.
- **Consejos prácticos**: preparación del perfil académico, certificaciones de idiomas, planificación financiera y uso de recursos oficiales.

Instrucciones de estilo y respuesta:
- Responde SIEMPRE en **español claro, profesional y motivador**.
- Organiza la respuesta en **secciones o pasos numerados** cuando expliques procesos.
- Usa **listas con viñetas** para resumir requisitos o documentos.
- Si la pregunta es muy amplia, primero ofrece un panorama general y luego invita al usuario a precisar más.
- Si la pregunta no tiene relación con becas o educación, responde brevemente indicando que tu especialidad son becas y redirige al tema.

Restricciones:
- Si no tienes información específica sobre una beca o convocatoria, dilo con transparencia y sugiere dónde puede buscar (páginas oficiales como Pronabec, embajadas, fundaciones, universidades).
- No inventes requisitos falsos; si no estás seguro, acláralo.
- Mantén siempre un tono empático, orientador y profesional.
"""

# =============== LLM y cadena ===============
llm = ChatOpenAI(
    model="gpt-4o",
    temperature=0.3,
    api_key=OPENAI_API_KEY,
)
prompt = ChatPromptTemplate.from_messages([
    ("system", system_prompt),
    ("human", "{pregunta}")
])
parser = StrOutputParser()
chain = prompt | llm | parser

# =============== Carga de documentos ===============
def load_pptx_text(path: str) -> str:
    """Extrae texto de PPTX usando python-pptx (sin dependencias pesadas)."""
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts).strip()

def cargar_documento(archivo_subido):
    """Carga y extrae texto de PDF, TXT, DOCX y PPTX."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(archivo_subido.name)[1]) as tmp_file:
            tmp_file.write(archivo_subido.getvalue())
            tmp_path = tmp_file.name

        extension = os.path.splitext(archivo_subido.name)[1].lower()

        if extension == '.pdf':
            loader = PyPDFLoader(tmp_path)
            documentos = loader.load()
            texto_completo = "\n\n".join([doc.page_content for doc in documentos])
        elif extension == '.txt':
            loader = TextLoader(tmp_path, encoding='utf-8')
            documentos = loader.load()
            texto_completo = "\n\n".join([doc.page_content for doc in documentos])
        elif extension in ['.docx']:
            loader = Docx2txtLoader(tmp_path)
            documentos = loader.load()
            texto_completo = "\n\n".join([doc.page_content for doc in documentos])
        elif extension in ['.pptx']:
            texto_completo = load_pptx_text(tmp_path)
        else:
            os.unlink(tmp_path)
            return None, f"❌ Formato no soportado: {extension}. Usa PDF, TXT, DOCX o PPTX."

        os.unlink(tmp_path)
        return texto_completo, None

    except Exception as e:
        return None, f"❌ Error al procesar el archivo: {str(e)}"

# =============== Sidebar: subir archivos ===============
st.sidebar.header("Subir Documentos")
st.sidebar.markdown("Puedes subir documentos relacionados con becas para que los analice:")

archivo_subido = st.sidebar.file_uploader(
    "Selecciona un archivo",
    type=['pdf', 'txt', 'docx', 'pptx'],
    help="Formatos soportados: PDF, TXT, DOCX, PPTX"
)

documento_texto = None
if archivo_subido is not None:
    with st.sidebar:
        with st.spinner(f"📊 Procesando {archivo_subido.name}..."):
            texto, error = cargar_documento(archivo_subido)
            if error:
                st.error(error)
            else:
                documento_texto = texto
                st.success(f"**{archivo_subido.name}** cargado")
                st.info(f"Tamaño: {len(texto):,} caracteres")

st.sidebar.markdown("---")
st.sidebar.info("""
**💡 Tipos de documentos útiles:**
- Convocatorias de becas
- Requisitos y bases
- Formularios de postulación
- Guías de aplicación
- Cartas de motivación
""")

# =============== Input del usuario ===============
user_input = st.chat_input("💭 Escribe tu pregunta sobre becas...")

if user_input:
    st.session_state.messages.append({"role": "user", "content": user_input})

    pregunta_final = user_input
    if documento_texto:
        texto_limitado = documento_texto[:8000]
        pregunta_final = f"""
Pregunta del usuario: {user_input}

Documento adjunto ({archivo_subido.name}):
{texto_limitado}

INSTRUCCIONES ESTRICTAS PARA EL ANÁLISIS:
1. Analiza el documento SOLO si está relacionado con becas para programas académicos de pregrado, maestría o doctorado
2. Si el documento trata sobre otros temas (cursos cortos, talleres, pasantías no académicas, etc.):
   - Detén el análisis inmediatamente
   - Responde: "El documento analizado no corresponde a becas para educación superior formal"
   - No proporciones ningún resumen o información del documento
3. Si el documento SÍ es sobre becas académicas:
   - Proporciona un resumen claro del contenido
   - Identifica requisitos, plazos, beneficios específicos
   - Ofrece recomendaciones basadas en el documento
   - Señala información faltante si aplica

Responde la pregunta del usuario basándote estrictamente en estas instrucciones.
"""

    with st.spinner("Analizando tu consulta..."):
        try:
            respuesta = chain.invoke({"pregunta": pregunta_final})
            st.session_state.messages.append({"role": "assistant", "content": respuesta})
            st.rerun()
        except Exception as e:
            error_msg = f"❌ Error: {str(e)}"
            st.session_state.messages.append({"role": "assistant", "content": error_msg})
            st.rerun()
